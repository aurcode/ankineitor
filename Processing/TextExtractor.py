from dotenv import load_dotenv
import os
import io
import re
import PyPDF2
from typing import List, Union, Dict
from pptx import Presentation
import pandas as pd
import jieba.posseg as pseg
from loguru import logger

# Load environment variables
load_dotenv()

import hashlib

'''class MongoHandler:
    def __init__(self, collection_name: str = 'file_results'):
        self.client = MongoDBClient()
        self.collection_name = collection_name
        self.field_name = "file_hash"

    def get_result(self, file_name: str) -> Union[None, pd.DataFrame]:
        """Retrieve the result from MongoDB if it already exists."""
        file_hash = self._get_file_hash(file_name)
        result = self.client.find_record(key=file_hash,
                                         collection_name=self.collection_name,
                                         field_name=self.field_name)
        if result:
            # Convert the stored result back to a DataFrame
            return pd.read_json(result['data'])
        return None

    def save_result(self, file_name: str, df: pd.DataFrame):
        """Save the processing results into MongoDB."""
        file_hash = self._get_file_hash(file_name)
        data = df.to_json()  # Save the DataFrame as JSON
        record = {self.field_name: file_hash, "data": data}

        try:
            self.client.insert_record(
                record=record,
                columns=["file_name", "data"],
                collection_name="results",
                field_name="file_name"
            )
            logger.info(f"Record for '{file_hash}' saved in MongoDB.")
        except Exception as e:
            logger.error(f"Error saving record for '{file_hash}' in MongoDB: {e}")

    def _get_file_hash(self, file_name: str) -> str:
        """Generate a unique hash for the file."""
        return hashlib.md5(file_name.encode()).hexdigest()
'''

class FileHandler:
    """Handles extraction of text from different file formats."""

    def __init__(self, file_content: bytes, test: bool = False):
        self.file_content = file_content
        self.test = test

    def extract_text(self, file_name: str) -> List[str]:
        """Dispatch to the correct method based on file extension."""
        file_extension = file_name.lower().split('.')[-1]

        if file_extension == 'pdf':
            return self._extract_text_from_pdf()
        elif file_extension == 'pptx':
            return self._extract_text_from_pptx()
        elif file_extension == 'txt':
            return self._extract_text_from_txt()
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    def _extract_text_from_txt(self) -> List[str]:
        """Extract text from a txt file."""
        text = self.file_content.decode('utf-8')
        return [text]

    def _extract_text_from_pdf(self) -> List[str]:
        """Extract text from a PDF file."""
        pdf_file = io.BytesIO(self.file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = []
        num_pages = len(pdf_reader.pages)

        # Optionally limit the pages for testing purposes
        max_page = 3 if self.test else num_pages
        for page_num in range(min(max_page, num_pages)):
            page = pdf_reader.pages[page_num]
            text.append(page.extract_text())

        return text

    def _extract_text_from_pptx(self) -> List[str]:
        """Extract text from a PowerPoint file."""
        pptx_file = io.BytesIO(self.file_content)
        presentation = Presentation(pptx_file)
        text = []

        max_page = 3 if self.test else len(presentation.slides)
        for slide in presentation.slides[:max_page]:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text.append(run.text)
        return text


class TextExtractor:
    """Class responsible for text extraction and processing from uploaded files."""

    def __init__(self, uploaded_files: Dict[str, bytes], dev_enabled: bool = os.getenv('DEBUG')):
        self.dev_enabled = dev_enabled
        self.uploaded_files = uploaded_files
        self.text: Union[str, List[str]] = []
        self.phrases: List[str] = []

    def extract_text(self):
        """Extract text from all uploaded files."""
        for file_name, file_content in self.uploaded_files.items():
            file_handler = FileHandler(file_content, self.dev_enabled)
            file_text = file_handler.extract_text(file_name)
            self.text.extend(file_text)
        self.text = '。'.join(self.text)

    def separated_chinese_characters(self, phrases=False) -> pd.DataFrame:
        """Separate and return Chinese characters along with their frequency and part-of-speech."""
        chinese_regex = r'[\u4E00-\u9FFF]+'
        chinese_characters = '。'.join(re.findall(chinese_regex, self.text))
        seg_list = list(pseg.cut(chinese_characters, use_paddle=True))

        # Create a DataFrame from the segmented list
        df = pd.DataFrame(set(seg_list), columns=["word", "part"])
        df = df[df['word'] != '。']
        df = df.groupby('word')['part'].agg(lambda x: ', '.join(x)).reset_index()

        df_aux = pd.DataFrame(seg_list, columns=["word", "part"])
        df_aux['frequency'] = df_aux['word'].map(df_aux['word'].value_counts())

        if phrases:
            self.phrases = self.extract_phrases()

        # Merge frequency data and return cleaned DataFrame
        return pd.merge(df, df_aux, on='word', how='left')[['word', 'part_x', 'frequency']] \
            .drop_duplicates(subset='word') \
            .rename(columns={'part_x': 'part'}) \
            .sort_values(by='frequency', ascending=False) \
            .reset_index(drop=True)

    def extract_phrases(self, split_n: int = 12, min_characters_n: int = 6) -> List[str]:
        """Extract meaningful phrases from text."""
        text_new = self._format_text(self.text, split_n)
        para = self._apply_text_formatting(text_new)
        return [i for i in para if len(i) >= min_characters_n and re.search(r'[\u4e00-\u9fff]+', i)]

    def _apply_text_formatting(self, text: str) -> List[str]:
        """Apply common text formatting (split by punctuation)."""
        para = re.sub(r"([。！？\?])([^”'])", r"\1\n\2", text)
        para = re.sub(r"(\.{6})([^”'])", r"\1\n\2", para)
        para = re.sub(r"(\…{2})([^”'])", r"\1\n\2", para)
        para = re.sub(r"([。！？\?][”'])([^，。！？\?])", r'\1\n\2', para)
        para = para.rstrip()
        para = list(set(para.split("\n")))
        return para

    def _format_text(self, text: str, split_n: int) -> str:
        """Format text into smaller segments."""
        lines = text.split('\n')
        result = []
        current_line = ''

        for line in lines:
            if len(current_line) > split_n:
                result.append(current_line + line)
                current_line = ''
            elif len(line) <= split_n:
                result.append(line + '\n')
            else:
                current_line = line

        if current_line:
            result.append(current_line)
        return ''.join(result)[:2000]

    @staticmethod
    def read_files_to_uploaded(file_paths: List[str]) -> Dict[str, bytes]:
        """Reads files from the given paths and returns a dictionary of file content."""
        uploaded_files = {}
        for file_path in file_paths:
            with open(file_path, 'rb') as file:
                uploaded_files[os.path.basename(file_path)] = file.read()
        return uploaded_files


if __name__ == '__main__':
    uploaded_files = TextExtractor.read_files_to_uploaded(['./电路CLASSES.pdf'])
    extractor = TextExtractor(uploaded_files, dev_enabled=False)
    extractor.extract_text()
    df = extractor.separated_chinese_characters()
    logger.info(f"Extracted {len(df)} rows of data.")
    print(df.head())